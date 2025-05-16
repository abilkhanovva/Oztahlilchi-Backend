import os
import docx
import PyPDF2
from pptx import Presentation
from rapidfuzz import process
from bs4 import BeautifulSoup

def sanitize_for_telegram(html_text: str) -> str:
    """
    HTML ichidagi <span class="error-word"> teglarini <b> teglariga almashtiradi.
    Telegram uchun matnni tayyorlashda ishlatiladi.
    """
    soup = BeautifulSoup(html_text, "html.parser")
    for span in soup.find_all("span", class_="error-word"):
        new_tag = soup.new_tag("b")
        new_tag.string = span.get_text()
        span.replace_with(new_tag)
    if soup.body:
        return str(soup.body.decode_contents())
    return str(soup)

def load_database():
    """
    So'zlar bazasini fayllardan yuklaydi va yagona to'plamga birlashtiradi.
    Yo'lni o'zingizga moslab o'zgartiring.
    """
    # Agar loyiha strukturasida bo'lsa, __file__ ga nisbatan:
    base_dir = os.path.dirname(os.path.abspath(__file__))
    # Fayllar asosiy papkada joylashgan bo'lsa, bir pog'ona yuqoriga ko'taring:
    latin_baza_path = os.path.join(base_dir, "..", "latin_baza.txt")
    fully_baza_path = os.path.join(base_dir, "..", "fully_baza.txt")

    latin = set()
    fully = set()

    with open(latin_baza_path, "r", encoding="utf-8") as f:
        latin.update(word.strip().lower() for word in f if word.strip())

    with open(fully_baza_path, "r", encoding="utf-8") as f:
        fully.update(word.strip().lower() for word in f if word.strip())

    return latin.union(fully)

# Umumiy so'zlar bazasi
DATABASE = load_database()

def analyze_text(text, max_suggestions=3, similarity_threshold=70, min_length_ratio=0.7):
    """
    Matndagi so'zlarni tekshiradi.
    To'g'ri so'zlar o'z holicha qoladi.
    Xato so'zlar <span> bilan belgilab, tavsiya so'zlar bilan birga qaytariladi.
    """
    words = text.strip().replace("\n", " ").split()
    corrected_words = []
    error_words = []
    suggestions_map = {}

    for word in words:
        clean_word = word.strip(".,?!;:()[]{}\"'«»").lower()

        if clean_word in DATABASE:
            corrected_words.append(word)
        else:
            error_words.append(clean_word)

            matches = process.extract(
                clean_word,
                DATABASE,
                limit=10,
                score_cutoff=similarity_threshold
            )

            filtered = [
                match[0]
                for match in matches
                if len(match[0]) >= len(clean_word) * min_length_ratio
            ][:max_suggestions]

            suggestions_map[clean_word] = filtered

            safe_suggestions = ",".join(filtered) if filtered else ""
            corrected_words.append(
                f'<span class="error-word" data-suggestions="{safe_suggestions}">{word}</span>'
            )

    corrected_text = " ".join(corrected_words)
    return corrected_text, error_words, suggestions_map

def extract_text_from_file(file):
    """
    Fayldan matn ajratib oladi.
    Qo'llab-quvvatlanadigan formatlar: txt, docx, pdf, pptx.
    """
    filename = file.filename
    ext = filename.split(".")[-1].lower()

    try:
        if ext == "txt":
            return file.read().decode("utf-8")

        elif ext == "docx":
            doc = docx.Document(file)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext == "pdf":
            reader = PyPDF2.PdfReader(file)
            texts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
            return "\n".join(texts)

        elif ext == "pptx":
            prs = Presentation(file)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)

    except Exception as e:
        print(f"[Xatolik] Faylni o‘qishda muammo: {e}")

    return ""
